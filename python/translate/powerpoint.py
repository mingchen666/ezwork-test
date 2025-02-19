import threading
import pptx
import translate
import common
import os
import sys
import time
import datetime

def start(trans):
    # 允许的最大线程
    threads=trans['threads']
    if threads is None or int(threads)<0:
        max_threads=10
    else:
        max_threads=int(threads)
    # 当前执行的索引位置
    run_index=0
    start_time = datetime.datetime.now()
    wb = pptx.Presentation(trans['file_path']) 
    slides = wb.slides
    texts=[]
    for slide in slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                text=paragraph.text
                if text!=None and len(text)>0 and not common.is_all_punc(text):
                    texts.append({"text":text, "complete":False})
    
    max_run=max_threads if len(texts)>max_threads else len(texts)
    before_active_count=threading.activeCount()
    event=threading.Event()
    while run_index<=len(texts)-1:
        if threading.activeCount()<max_run+before_active_count:
            if not event.is_set():
                thread = threading.Thread(target=translate.get,args=(trans,event,texts,run_index))
                thread.start()
                run_index+=1
            else:
                return False
    
    while True:
        complete=True
        for text in texts:
            if not text['complete']:
                complete=False
        if complete:
            break
        else:
            time.sleep(1)

    text_count=0
    for slide in slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                text=paragraph.text
                if text!=None and len(text)>0 and not common.is_all_punc(text) and len(texts)>0:
                    item=texts.pop(0)
                    paragraph.text=item['text']
                    text_count+=item['count']

    wb.save(trans['target_file'])
    end_time = datetime.datetime.now()
    spend_time=common.display_spend(start_time, end_time)
    translate.complete(trans,text_count,spend_time)
    return True


